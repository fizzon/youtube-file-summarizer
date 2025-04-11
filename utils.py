# utils.py

import textwrap
import streamlit as st # Потрібен для st.error
import logging
import io # Для роботи з байтовими потоками завантажених файлів

# Бібліотеки для обробки файлів (переконайтесь, що вони встановлені)
try:
    import PyPDF2
except ImportError:
    st.error("Будь ласка, встановіть бібліотеку pypdf2: pip install pypdf2")
    # Можна зупинити виконання, якщо бібліотека критична
    # st.stop()
try:
    from pptx import Presentation
except ImportError:
    st.error("Будь ласка, встановіть бібліотеку python-pptx: pip install python-pptx")
    # st.stop()


# --- Функція розбиття тексту (існуюча) ---
def split_text(text, max_length=3000):
    """Розбиває текст на частини заданої максимальної довжини."""
    if not isinstance(text, str):
        logging.error(f"split_text отримав не текстовий тип: {type(text)}")
        return [] # Повертаємо порожній список у разі помилки типу
    # Замінюємо нерозривні пробіли на звичайні перед розбиттям
    text = text.replace('\xa0', ' ')
    return textwrap.wrap(text, width=max_length, break_long_words=False, break_on_hyphens=False)

# --- НОВІ Функції для вилучення тексту з файлів ---

def extract_text_from_pdf(uploaded_file):
    """Вилучає текст з PDF-файлу."""
    text = ""
    try:
        # Читаємо завантажений файл як байтовий потік
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(uploaded_file.getvalue()))
        logging.info(f"Обробка PDF: {uploaded_file.name}, кількість сторінок: {len(pdf_reader.pages)}")
        for i, page in enumerate(pdf_reader.pages):
            page_text = page.extract_text()
            if page_text:
                 text += page_text + "\n"
            else:
                 logging.warning(f"Не вдалося вилучити текст зі сторінки {i+1} файлу {uploaded_file.name}")
        logging.info(f"Текст з PDF {uploaded_file.name} вилучено успішно.")
        return text.strip() # Повертаємо текст без зайвих пробілів на початку/кінці
    except Exception as e:
        st.error(f"Помилка обробки PDF '{uploaded_file.name}': {e}")
        logging.error(f"Error processing PDF {uploaded_file.name}: {e}")
        return None # Повертаємо None у разі помилки

def extract_text_from_pptx(uploaded_file):
    """Вилучає текст з PPTX-файлу."""
    text = ""
    try:
        # Читаємо завантажений файл як байтовий потік
        prs = Presentation(io.BytesIO(uploaded_file.getvalue()))
        logging.info(f"Обробка PPTX: {uploaded_file.name}, кількість слайдів: {len(prs.slides)}")
        for slide_num, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    # Додаємо текст з фігури, видаляючи зайві пробіли
                    shape_text = shape.text.strip()
                    if shape_text: # Додаємо тільки якщо текст не порожній
                         slide_text += shape_text + "\n" # Новий рядок після тексту кожної фігури
            if slide_text: # Додаємо текст слайду, якщо він не порожній
                 text += f"--- Слайд {slide_num + 1} ---\n{slide_text.strip()}\n\n"
            else:
                 logging.warning(f"Не знайдено тексту на слайді {slide_num + 1} файлу {uploaded_file.name}")

        logging.info(f"Текст з PPTX {uploaded_file.name} вилучено успішно.")
        return text.strip() # Повертаємо текст без зайвих пробілів на початку/кінці
    except Exception as e:
        st.error(f"Помилка обробки PPTX '{uploaded_file.name}': {e}")
        logging.error(f"Error processing PPTX {uploaded_file.name}: {e}")
        return None # Повертаємо None у разі помилки

def extract_text_from_txt(uploaded_file):
     """Вилучає текст з TXT-файлу."""
     try:
          # Читаємо байтовий потік і декодуємо як UTF-8
          text = uploaded_file.getvalue().decode("utf-8")
          logging.info(f"Текст з TXT {uploaded_file.name} вилучено успішно.")
          return text.strip()
     except UnicodeDecodeError:
          st.error(f"Помилка декодування TXT файлу '{uploaded_file.name}'. Спробуйте зберегти файл у кодуванні UTF-8.")
          logging.error(f"UnicodeDecodeError processing TXT {uploaded_file.name}")
          return None
     except Exception as e:
          st.error(f"Помилка обробки TXT '{uploaded_file.name}': {e}")
          logging.error(f"Error processing TXT {uploaded_file.name}: {e}")
          return None # Повертаємо None у разі помилки
